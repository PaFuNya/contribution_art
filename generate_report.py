"""
生成项目报告 (docx) 和 汇报PPT (pptx)
"""
import sys
import json
import os
sys.stdout.reconfigure(encoding='utf-8')

import pandas as pd
import numpy as np
from pathlib import Path

# ==================== 报告生成 ====================

def generate_docx_report():
    """生成 Word 报告"""
    from docx import Document
    from docx.shared import Inches, Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT

    doc = Document()

    # 样式设置
    style = doc.styles['Normal']
    style.font.name = 'SimSun'
    style.font.size = Pt(12)
    style.paragraph_format.line_spacing = 1.5

    # === 封面 ===
    for _ in range(6):
        doc.add_paragraph()

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run('动漫数据分析报告')
    run.font.size = Pt(28)
    run.bold = True

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run('—— 基于 MyAnimeList 数据集的完整数据分析')
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(100, 100, 100)

    info = doc.add_paragraph()
    info.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = info.add_run('\n\n《数据分析应用》期末综合任务\n\n')
    run.font.size = Pt(14)

    doc.add_page_break()

    # === 目录 ===
    doc.add_heading('目录', level=1)
    toc_items = [
        '一、研究背景与问题定义',
        '二、数据来源与字段说明',
        '三、数据清洗与预处理',
        '四、探索性数据分析（EDA）',
        '五、可视化分析',
        '六、特征工程',
        '七、建模方法与参数说明',
        '八、模型评估结果',
        '九、大模型辅助分析',
        '十、聚类分析',
        '十一、项目结论与建议',
        '十二、项目不足与改进方向',
    ]
    for item in toc_items:
        doc.add_paragraph(item, style='List Number')

    doc.add_page_break()

    # === 一、研究背景 ===
    doc.add_heading('一、研究背景与问题定义', level=1)
    doc.add_paragraph(
        '随着全球动漫产业的快速发展，动漫作品数量逐年增长。MyAnimeList 作为全球最大的动漫社区平台，'
        '积累了海量的用户评分、收藏、关注等行为数据。本项目基于该平台的动漫数据集，旨在回答以下核心问题：'
    )
    questions = [
        '动漫评分的分布规律是什么？哪些因素最显著地影响评分？',
        '能否基于动漫的类型、集数、工作室等特征预测其评分？',
        '能否将动漫自动分群？各群体具有什么特征？',
    ]
    for q in questions:
        doc.add_paragraph(q, style='List Bullet')

    # === 二、数据来源 ===
    doc.add_heading('二、数据来源与字段说明', level=1)
    doc.add_paragraph(
        '数据来源：MyAnimeList 平台公开数据\n'
        '数据规模：15,000 行 × 20 列\n'
        '时间跨度：1961 - 2025 年'
    )

    # 字段表
    table = doc.add_table(rows=13, cols=3)
    table.style = 'Light Grid Accent 1'
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    headers = ['字段名', '数据类型', '说明']
    for i, h in enumerate(headers):
        table.rows[0].cells[i].text = h

    fields = [
        ('番ID', 'int64', '动漫唯一标识'),
        ('番名', 'object', '动漫名称'),
        ('分数', 'float64', '用户综合评分（1-10）'),
        ('类型', 'object', '动漫类型（多值，逗号分隔）'),
        ('主题', 'object', '主题标签'),
        ('受众群体', 'object', '少年/青年/少女等'),
        ('视频类型', 'object', '电视/电影/影碟/ONA等'),
        ('集数', 'float64', '总集数'),
        ('年份', 'float64', '首播年份'),
        ('工作室', 'object', '制作公司'),
        ('源材料', 'object', '漫画/原创/游戏/轻小说等'),
        ('收藏/关注人数/社区成员数', 'int64', '人气指标'),
    ]
    for i, (name, dtype, desc) in enumerate(fields):
        table.rows[i + 1].cells[0].text = name
        table.rows[i + 1].cells[1].text = dtype
        table.rows[i + 1].cells[2].text = desc

    # === 三、数据清洗 ===
    doc.add_heading('三、数据清洗与预处理', level=1)
    doc.add_paragraph('数据清洗流程如下：')

    clean_steps = [
        '时长字段转换：将"24分钟"格式转为数值（分钟）',
        '多值列拆分：将"类型"和"主题"字段拆分为列表',
        '缺失值处理：集数用中位数填充，分类列用"Unknown"填充',
        '高缺失列处理：删除缺失率 > 70% 的"季节"和"年份和季节"列',
        '删除关键字段缺失行：年份和类型缺失的行（约10,000条）',
        '类型转换：年份转整数，集数转整数',
    ]
    for step in clean_steps:
        doc.add_paragraph(step, style='List Bullet')

    doc.add_paragraph(f'清洗后数据量：4,562 行 × 22 列')

    # === 四、EDA ===
    doc.add_heading('四、探索性数据分析（EDA）', level=1)

    doc.add_heading('4.1 评分分布', level=2)
    doc.add_paragraph(
        '动漫评分呈近似正态分布，均值为 6.95，标准差为 0.71。'
        '中位数为 6.93，说明评分分布较为对称。'
        '评分集中在 6.0-8.0 区间，极端高分（>9.0）和极端低分（<5.0）均较少。'
    )
    if Path('output/dist_分数.png').exists():
        doc.add_picture('output/dist_分数.png', width=Inches(6))
        doc.add_paragraph('图1：动漫评分分布').alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading('4.2 类型分布', level=2)
    doc.add_paragraph(
        '最常见的动漫类型为：喜剧（1157部）、动作（353部）、奇幻（184部）、冒险（205部）、幻想（417部）。'
        '很多动漫属于多种类型组合，如"动作,冒险,奇幻"。'
    )
    if Path('output/bar_genre_top15.png').exists():
        doc.add_picture('output/bar_genre_top15.png', width=Inches(5.5))
        doc.add_paragraph('图2：Top 15 动漫类型').alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading('4.3 年份趋势', level=2)
    doc.add_paragraph(
        '2010年后动漫产出显著增加，2024年达到峰值（216部）。'
        '平均评分在不同年份间波动较小，整体保持稳定。'
    )
    if Path('output/bar_yearly_count.png').exists():
        doc.add_picture('output/bar_yearly_count.png', width=Inches(6))
        doc.add_paragraph('图3：每年动漫产出数量').alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading('4.4 相关性分析', level=2)
    doc.add_paragraph(
        '数值特征相关性热力图显示，收藏数、关注人数、社区成员数之间高度正相关。'
        '评分与人气指标（收藏、关注）呈中等正相关。'
    )
    if Path('output/correlation_heatmap.png').exists():
        doc.add_picture('output/correlation_heatmap.png', width=Inches(5))
        doc.add_paragraph('图4：数值特征相关性热力图').alignment = WD_ALIGN_PARAGRAPH.CENTER

    # === 五、可视化 ===
    doc.add_heading('五、可视化分析', level=1)

    doc.add_heading('5.1 视频类型对比', level=2)
    doc.add_paragraph(
        '电视动画和电影的评分分布较为集中，ONA（网络动画）的评分方差较大。'
    )
    if Path('output/box_视频类型_vs_分数.png').exists():
        doc.add_picture('output/box_视频类型_vs_分数.png', width=Inches(5.5))
        doc.add_paragraph('图5：不同视频类型的评分对比').alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading('5.2 源材料对比', level=2)
    doc.add_paragraph(
        '漫画改编和原创作品是两大主要来源。视觉小说改编作品评分方差最大。'
    )
    if Path('output/box_源材料_vs_分数.png').exists():
        doc.add_picture('output/box_源材料_vs_分数.png', width=Inches(5.5))
        doc.add_paragraph('图6：不同源材料的评分对比').alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading('5.3 工作室对比', level=2)
    doc.add_paragraph(
        'ToeiAnimation（700部）、Sunrise（513部）、J.C.Staff（359部）是产量最高的三家工作室。'
        'KyotoAnimation 虽然产量不高（116部），但评分中位数较高。'
    )
    if Path('output/box_studio_vs_score.png').exists():
        doc.add_picture('output/box_studio_vs_score.png', width=Inches(6))
        doc.add_paragraph('图7：Top 15 工作室评分分布').alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading('5.4 人气与评分关系', level=2)
    doc.add_paragraph(
        '收藏数与评分呈正相关，但存在明显的长尾效应：少数作品获得极高收藏。'
    )
    if Path('output/scatter_收藏_vs_分数.png').exists():
        doc.add_picture('output/scatter_收藏_vs_分数.png', width=Inches(5.5))
        doc.add_paragraph('图8：收藏数 vs 评分').alignment = WD_ALIGN_PARAGRAPH.CENTER

    # === 六、特征工程 ===
    doc.add_heading('六、特征工程', level=1)
    doc.add_paragraph('最终构建了 28 个特征，分为三类：')

    table2 = doc.add_table(rows=4, cols=3)
    table2.style = 'Light Grid Accent 1'
    table2.rows[0].cells[0].text = '特征类型'
    table2.rows[0].cells[1].text = '数量'
    table2.rows[0].cells[2].text = '说明'
    table2.rows[1].cells[0].text = '数值特征'
    table2.rows[1].cells[1].text = '9'
    table2.rows[1].cells[2].text = '集数、时长、人气排名、收藏、关注人数等'
    table2.rows[2].cells[0].text = '编码特征'
    table2.rows[2].cells[1].text = '4'
    table2.rows[2].cells[2].text = '视频类型、源材料、受众群体、年龄适度'
    table2.rows[3].cells[0].text = '类型特征'
    table2.rows[3].cells[1].text = '15'
    table2.rows[3].cells[2].text = 'Top 15 类型的 one-hot 编码'

    doc.add_paragraph('\n特征工程步骤：')
    feat_steps = [
        '人气综合得分：收藏、关注、社区成员归一化后加权求和',
        '时间特征：番龄（2025-年份）、是否经典（≥10年）、年代分箱',
        '类型 one-hot：Top 15 类型各生成一个 0/1 特征',
        '分类编码：Label Encoding 处理视频类型、源材料等',
    ]
    for s in feat_steps:
        doc.add_paragraph(s, style='List Bullet')

    # === 七、建模 ===
    doc.add_heading('七、建模方法与参数说明', level=1)
    doc.add_paragraph(
        '本项目采用回归任务预测动漫评分，使用 80/20 划分训练集/测试集，random_state=42。'
    )

    table3 = doc.add_table(rows=6, cols=4)
    table3.style = 'Light Grid Accent 1'
    for i, h in enumerate(['模型', '关键参数', '特点']):
        table3.rows[0].cells[i].text = h
    models_info = [
        ('LinearRegression', '默认', '线性基线'),
        ('Ridge', 'alpha=1.0', 'L2正则化'),
        ('Lasso', 'alpha=0.01', 'L1正则化，特征选择'),
        ('RandomForest', 'n_estimators=100, max_depth=10', '集成学习，非线性'),
        ('GradientBoosting', 'n_estimators=100, max_depth=5, lr=0.1', '梯度提升，强泛化'),
    ]
    for i, (name, params, desc) in enumerate(models_info):
        table3.rows[i + 1].cells[0].text = name
        table3.rows[i + 1].cells[1].text = params
        table3.rows[i + 1].cells[2].text = desc

    # === 八、模型评估 ===
    doc.add_heading('八、模型评估结果', level=1)

    table4 = doc.add_table(rows=6, cols=4)
    table4.style = 'Light Grid Accent 1'
    for i, h in enumerate(['模型', 'RMSE', 'MAE', 'R²']):
        table4.rows[0].cells[i].text = h
    eval_data = [
        ('LinearRegression', '0.4946', '0.3905', '0.5485'),
        ('Ridge', '0.4945', '0.3905', '0.5486'),
        ('Lasso', '0.5007', '0.3968', '0.5372'),
        ('RandomForest', '0.3964', '0.3003', '0.7100'),
        ('GradientBoosting', '0.3863', '0.2923', '0.7246'),
    ]
    for i, (name, rmse, mae, r2) in enumerate(eval_data):
        table4.rows[i + 1].cells[0].text = name
        table4.rows[i + 1].cells[1].text = rmse
        table4.rows[i + 1].cells[2].text = mae
        table4.rows[i + 1].cells[3].text = r2

    doc.add_paragraph(
        '\nGradientBoosting 模型表现最佳，R² = 0.7246，能解释约 72% 的评分变异。'
    )

    if Path('output/feature_importance.png').exists():
        doc.add_picture('output/feature_importance.png', width=Inches(5.5))
        doc.add_paragraph('图9：特征重要性 Top 15').alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph(
        '特征重要性分析显示，收藏数是预测评分的最强特征（重要性 69%），'
        '其次是集数（5.4%）、番龄（4.4%）、关注人数（3.9%）。'
    )

    if Path('output/regression_pred_vs_actual.png').exists():
        doc.add_picture('output/regression_pred_vs_actual.png', width=Inches(4.5))
        doc.add_paragraph('图10：预测值 vs 实际值').alignment = WD_ALIGN_PARAGRAPH.CENTER

    # === 九、大模型辅助分析 ===
    doc.add_heading('九、大模型辅助分析', level=1)
    doc.add_paragraph(
        '本项目使用阿里百炼 DashScope API（模型：kimi-k2.6）对分析结果进行结构化总结。'
        '分别测试了 temperature=0.2、0.7、1.2 三个参数，对比输出差异。'
    )

    # 加载 LLM 结果
    try:
        with open('output/llm_analysis.json', 'r', encoding='utf-8') as f:
            llm_result = json.load(f)

        doc.add_heading('LLM 分析结论', level=2)
        doc.add_paragraph(f"项目概述：{llm_result.get('项目概述', '')}")

        doc.add_heading('关键发现', level=3)
        for finding in llm_result.get('关键发现', []):
            doc.add_paragraph(finding, style='List Bullet')

        doc.add_heading('模型评价', level=3)
        doc.add_paragraph(llm_result.get('模型评价', ''))

        doc.add_heading('聚类解读', level=3)
        doc.add_paragraph(llm_result.get('聚类解读', ''))

        doc.add_heading('业务建议', level=3)
        for advice in llm_result.get('业务建议', []):
            doc.add_paragraph(advice, style='List Bullet')

        doc.add_heading('局限性', level=3)
        for lim in llm_result.get('局限性', []):
            doc.add_paragraph(lim, style='List Bullet')
    except Exception as e:
        doc.add_paragraph(f'（LLM 分析结果加载失败：{e}）')

    doc.add_heading('Temperature 对比分析', level=2)
    doc.add_paragraph(
        'temperature=0.2：输出最客观、条理清晰，适合数据分析报告场景。\n'
        'temperature=0.7：有一定灵活性，建议更丰富，兼顾客观性。\n'
        'temperature=1.2：输出更随机，可能出现偏离数据的表述。\n'
        '结论：数据分析场景推荐 temperature=0.2，输出稳定可靠。'
    )

    # === 十、聚类 ===
    doc.add_heading('十、聚类分析', level=1)
    doc.add_paragraph(
        '使用 K-Means 聚类（K=4），通过轮廓系数确定最优聚类数。'
        'PCA 降维可视化显示各聚类有较好分离。'
    )

    table5 = doc.add_table(rows=5, cols=5)
    table5.style = 'Light Grid Accent 1'
    for i, h in enumerate(['聚类', '数量', '平均评分', '平均收藏', '特征描述']):
        table5.rows[0].cells[i].text = h
    cluster_desc = [
        ('聚类0', '1035', '6.50', '38', '长尾基础作品，低人气低评分'),
        ('聚类1', '1912', '6.94', '1242', '普通流行作品，市场腰部'),
        ('聚类2', '1484', '7.13', '1460', '主流优质作品，口碑较好'),
        ('聚类3', '131', '8.16', '44903', '现象级爆款，超高人气'),
    ]
    for i, (cl, cnt, score, fav, desc) in enumerate(cluster_desc):
        table5.rows[i + 1].cells[0].text = cl
        table5.rows[i + 1].cells[1].text = cnt
        table5.rows[i + 1].cells[2].text = score
        table5.rows[i + 1].cells[3].text = fav
        table5.rows[i + 1].cells[4].text = desc

    if Path('output/pca_clusters.png').exists():
        doc.add_picture('output/pca_clusters.png', width=Inches(5))
        doc.add_paragraph('图11：聚类 PCA 可视化').alignment = WD_ALIGN_PARAGRAPH.CENTER

    if Path('output/cluster_comparison.png').exists():
        doc.add_picture('output/cluster_comparison.png', width=Inches(5.5))
        doc.add_paragraph('图12：各聚类特征对比').alignment = WD_ALIGN_PARAGRAPH.CENTER

    # === 十一、结论 ===
    doc.add_heading('十一、项目结论与建议', level=1)

    doc.add_heading('核心结论', level=2)
    conclusions = [
        '收藏数是预测动漫评分的最强特征，说明用户情感投入与作品质量高度相关。',
        'GradientBoosting 模型能解释 72% 的评分变异，预测效果良好。',
        '动漫市场呈典型金字塔结构：2.9%头部作品、36%腰部作品、22.7%长尾作品。',
        '喜剧、动作、奇幻是最主流的类型组合。',
        'ToeiAnimation、Sunrise 是产量最高的工作室，但产量与评分并非正相关。',
    ]
    for c in conclusions:
        doc.add_paragraph(c, style='List Bullet')

    doc.add_heading('建议', level=2)
    suggestions = [
        '新番制作应重点关注收藏数的提升，这是评分的最强信号。',
        '漫画改编和原创是最稳定的题材来源，可优先考虑。',
        '运营策略应分层化：头部作品重点IP开发，腰部作品强化社区运营。',
    ]
    for s in suggestions:
        doc.add_paragraph(s, style='List Bullet')

    # === 十二、不足 ===
    doc.add_heading('十二、项目不足与改进方向', level=1)
    limitations = [
        '数据缺失：受众群体缺失70%、年份缺失68%，影响分析完整性。',
        '特征局限：未纳入导演、编剧、声优、制作经费等关键变量。',
        '模型局限：R²=0.72 仍有提升空间，可尝试深度学习或更多特征。',
        '因果性：收藏数与评分可能存在因果倒置（高评分带来更多收藏）。',
    ]
    for lim in limitations:
        doc.add_paragraph(lim, style='List Bullet')

    # 保存
    output_path = 'report/动漫数据分析报告.docx'
    Path('report').mkdir(exist_ok=True)
    doc.save(output_path)
    print(f"✅ 报告已生成: {output_path}")
    return output_path


# ==================== PPT 生成 ====================

def generate_pptx():
    """生成汇报 PPT"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_slide(title_text, content_items=None, image_path=None, table_data=None):
        """添加一页幻灯片"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 标题栏
        from pptx.util import Inches, Pt
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(44, 62, 80)

        # 分隔线
        from pptx.util import Emu
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12.3), Emu(28000))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(52, 152, 219)
        line.line.fill.background()

        y_offset = 1.5

        # 文字内容
        if content_items:
            textbox = slide.shapes.add_textbox(Inches(0.8), Inches(y_offset), Inches(11.5), Inches(5))
            tf = textbox.text_frame
            tf.word_wrap = True
            for i, item in enumerate(content_items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(52, 73, 94)
                p.space_after = Pt(8)
            y_offset += len(content_items) * 0.4 + 0.3

        # 图片
        if image_path and Path(image_path).exists():
            img_left = Inches(1.5) if not content_items else Inches(5)
            img_top = Inches(y_offset)
            img_width = Inches(10) if not content_items else Inches(7.5)
            pic = slide.shapes.add_picture(image_path, img_left, img_top, img_width)
            # 限制高度
            if pic.height > Inches(5):
                pic.height = Inches(5)

        # 表格
        if table_data:
            rows = len(table_data)
            cols = len(table_data[0])
            tbl = slide.shapes.add_table(rows, cols,
                                          Inches(1), Inches(y_offset),
                                          Inches(11), Inches(rows * 0.45)).table
            for r in range(rows):
                for c in range(cols):
                    cell = tbl.cell(r, c)
                    cell.text = str(table_data[r][c])
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(14)
                        if r == 0:
                            paragraph.font.bold = True

        return slide

    # === PPT 内容 ===

    # 1. 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = '🎬 动漫数据分析'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(44, 62, 80)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = '基于 MyAnimeList 数据集的完整数据分析'
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(100, 100, 100)
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = '\n《数据分析应用》期末综合任务'
    p3.font.size = Pt(20)
    p3.alignment = PP_ALIGN.CENTER

    # 2. 目录
    add_slide('目录', [
        '1. 研究背景与问题定义',
        '2. 数据来源与字段说明',
        '3. 数据清洗与预处理',
        '4. 探索性数据分析（EDA）',
        '5. 特征工程与建模',
        '6. 模型评估结果',
        '7. 聚类分析',
        '8. 大模型辅助分析',
        '9. 结论与建议',
    ])

    # 3. 研究背景
    add_slide('一、研究背景与问题定义', [
        '📌 研究背景：',
        '  • 全球动漫产业快速发展，MyAnimeList 积累了海量用户行为数据',
        '  • 数据驱动的分析可以揭示动漫市场的规律和趋势',
        '',
        '🎯 核心问题：',
        '  1. 哪些因素最显著地影响动漫评分？',
        '  2. 能否基于特征预测评分？',
        '  3. 能否将动漫自动分群？各群体有什么特征？',
    ])

    # 4. 数据来源
    add_slide('二、数据来源与字段说明', [
        '📊 数据概况：',
        '  • 来源：MyAnimeList 平台公开数据',
        '  • 规模：15,000 行 × 20 列',
        '  • 时间跨度：1961 - 2025 年',
        '',
        '📋 核心字段：',
        '  番名、分数、类型、主题、受众群体、视频类型、集数、',
        '  年份、工作室、源材料、时长、人气排名、收藏、关注人数、社区成员数',
    ])

    # 5. 数据清洗
    add_slide('三、数据清洗与预处理', [
        '🔧 清洗步骤：',
        '  1. 时长字段转换："24分钟" → 24（数值）',
        '  2. 多值列拆分：类型、主题拆分为列表',
        '  3. 缺失值处理：集数用中位数填充，分类列用"Unknown"',
        '  4. 删除高缺失列：季节（68.8%空）、年份和季节',
        '  5. 删除关键字段缺失行',
        '',
        '✅ 清洗结果：4,562 行 × 22 列（原 15,000 行）',
    ])

    # 6. EDA - 评分分布
    add_slide('四、EDA — 评分分布', [
        '📊 评分分布特征：',
        '  • 均值：6.95，标准差：0.71',
        '  • 集中在 6.0-8.0 区间',
        '  • 极端高分（>9.0）和低分（<5.0）较少',
    ], image_path='output/dist_分数.png')

    # 7. EDA - 类型分布
    add_slide('四、EDA — 类型分布', [
        '📊 Top 5 类型：',
        '  • 喜剧（1157部）→ 动作（353部）→ 幻想（417部）',
        '  • 冒险（205部）→ 奇幻（184部）',
        '',
        '💡 很多动漫是多种类型组合',
    ], image_path='output/bar_genre_top15.png')

    # 8. EDA - 年份趋势 + 相关性
    add_slide('四、EDA — 年份趋势与相关性', [
        '📈 年份趋势：2010年后产出显著增加',
        '🔥 相关性：收藏数、关注人数、社区成员数高度正相关',
    ], image_path='output/correlation_heatmap.png')

    # 9. 可视化 - 视频类型 + 源材料
    add_slide('五、可视化 — 视频类型与源材料', [
        '📺 视频类型：电视动画和电影评分分布集中',
        '📖 源材料：漫画改编和原创是两大主要来源',
    ], image_path='output/box_视频类型_vs_分数.png')

    # 10. 特征工程
    add_slide('六、特征工程', [
        '🔧 特征构建（共28个特征）：',
        '',
        '  数值特征（9个）：集数、时长、人气排名、收藏、关注人数等',
        '  编码特征（4个）：视频类型、源材料、受众群体、年龄适度',
        '  类型特征（15个）：Top 15 类型的 one-hot 编码',
        '',
        '  ➕ 人气综合得分（加权归一化）',
        '  ➕ 番龄、是否经典、年代分箱',
    ])

    # 11. 建模结果
    add_slide('七、建模结果', [
        '🏆 模型对比：',
    ], table_data=[
        ['模型', 'RMSE', 'MAE', 'R²'],
        ['LinearRegression', '0.4946', '0.3905', '0.5485'],
        ['Ridge', '0.4945', '0.3905', '0.5486'],
        ['Lasso', '0.5007', '0.3968', '0.5372'],
        ['RandomForest', '0.3964', '0.3003', '0.7100'],
        ['GradientBoosting', '0.3863', '0.2923', '0.7246'],
    ])

    # 12. 特征重要性
    add_slide('七、特征重要性', [
        '🔑 最重要特征：收藏数（69%）',
        '  • 其次：集数（5.4%）、番龄（4.4%）、关注人数（3.9%）',
        '  • 说明：用户行为数据比内容属性更能预测评分',
    ], image_path='output/feature_importance.png')

    # 13. 聚类分析
    add_slide('八、聚类分析', [
        '🔬 K-Means 聚类（K=4）：',
    ], table_data=[
        ['聚类', '数量', '平均评分', '平均收藏', '特征'],
        ['聚类0', '1035', '6.50', '38', '长尾基础作品'],
        ['聚类1', '1912', '6.94', '1242', '普通流行作品'],
        ['聚类2', '1484', '7.13', '1460', '主流优质作品'],
        ['聚类3', '131', '8.16', '44903', '现象级爆款'],
    ])

    # 14. 聚类可视化
    add_slide('八、聚类 PCA 可视化', [
        '💡 市场呈金字塔结构：2.9%头部 → 36%腰部 → 22.7%长尾',
    ], image_path='output/pca_clusters.png')

    # 15. 大模型辅助分析
    add_slide('九、大模型辅助分析', [
        '🤖 使用阿里百炼 DashScope API（kimi-k2.6）',
        '',
        '📋 分析内容：',
        '  • 对分析结果生成结构化总结',
        '  • 多 temperature 对比（0.2 / 0.7 / 1.2）',
        '',
        '💡 结论：',
        '  • temperature=0.2 最客观，适合数据分析报告',
        '  • temperature=1.2 输出更随机，可能偏离数据',
        '',
        '✅ 详见 output/llm_analysis.json 和 llm_multi_temperature.json',
    ])

    # 16. 结论
    add_slide('十、核心结论', [
        '✅ 1. 收藏数是预测评分的最强特征（重要性69%）',
        '✅ 2. GradientBoosting 模型 R²=0.72，能解释72%评分变异',
        '✅ 3. 动漫市场呈金字塔结构：头部2.9%、腰部36%、长尾22.7%',
        '✅ 4. 喜剧、动作、奇幻是最主流类型组合',
        '✅ 5. ToeiAnimation、Sunrise 产量最高，但产量≠评分',
    ])

    # 17. 不足与改进
    add_slide('十一、不足与改进', [
        '⚠️ 局限性：',
        '  • 数据缺失：受众群体70%、年份68%',
        '  • 特征局限：未纳入导演、编剧、声优等',
        '  • 模型 R²=0.72 仍有提升空间',
        '',
        '🔮 改进方向：',
        '  • 补充外部数据（制作经费、播出平台等）',
        '  • 尝试深度学习模型',
        '  • 引入文本分析（评论情感、简介关键词）',
    ])

    # 18. 致谢
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = '感谢聆听！'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(44, 62, 80)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = '\nQ & A'
    p2.font.size = Pt(32)
    p2.font.color.rgb = RGBColor(52, 152, 219)
    p2.alignment = PP_ALIGN.CENTER

    # 保存
    output_path = 'report/动漫数据分析汇报.pptx'
    Path('report').mkdir(exist_ok=True)
    prs.save(output_path)
    print(f"✅ PPT 已生成: {output_path}")
    return output_path


if __name__ == '__main__':
    print("=" * 60)
    print("生成报告和PPT")
    print("=" * 60)
    generate_docx_report()
    generate_pptx()
    print("\n✅ 全部生成完毕！文件在 report/ 目录下。")
